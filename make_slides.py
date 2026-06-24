"""
Generate slides.pptx — a presentation of the synthetic-data pipeline.

Run:  python make_slides.py
Embeds the rendered figures in reports/ and the verified benchmark numbers.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

REPORTS = Path("reports")

# ── palette ───────────────────────────────────────────────────────────────
INK    = RGBColor(0x1F, 0x2A, 0x37)   # near-black text
MUTED  = RGBColor(0x5B, 0x66, 0x70)
ACCENT = RGBColor(0x4E, 0x79, 0xA7)   # M1 blue
BG     = RGBColor(0xFF, 0xFF, 0xFF)
PANEL  = RGBColor(0xF3, 0xF5, 0xF7)
GOOD   = RGBColor(0x2E, 0x7D, 0x32)
WARN   = RGBColor(0xC6, 0x28, 0x28)
METHOD = {"M1": RGBColor(0x4E,0x79,0xA7), "M2": RGBColor(0xF2,0x8E,0x2B),
          "M3": RGBColor(0x59,0xA1,0x4F), "M4": RGBColor(0xE1,0x57,0x59),
          "M5": RGBColor(0xB0,0x7A,0xA1)}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid(); bg.fore_color.rgb = BG
    return s


def box(s, x, y, w, h):
    return s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)).text_frame


def _set(p, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, italic=False):
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Calibri"


def accent_bar(s, color=ACCENT):
    bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.22), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()


def header(s, title, kicker=None, color=ACCENT):
    accent_bar(s, color)
    tf = box(s, 0.6, 0.42, 12.2, 1.1)
    if kicker:
        _set(tf.paragraphs[0], kicker.upper(), 13, color, bold=True)
        p = tf.add_paragraph(); _set(p, title, 32, INK, bold=True)
    else:
        _set(tf.paragraphs[0], title, 32, INK, bold=True)


def bullets(s, items, x=0.85, y=1.9, w=11.6, h=5.0, size=18, gap=10):
    h = min(h, 7.2 - y)   # never run past the slide bottom
    tf = box(s, x, y, w, h); tf.word_wrap = True
    for i, it in enumerate(items):
        lvl = 0
        if isinstance(it, tuple):
            it, lvl = it
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = "•  " if lvl == 0 else "–  "
        _set(p, bullet + it, size - lvl*2, INK if lvl == 0 else MUTED, bold=(lvl == 0 and it.endswith(":")))
        p.space_after = Pt(gap)
        p.level = lvl


def fit_image(s, path, x, y, w, h):
    """Place an image fit (contain) within the box (x,y,w,h) in inches, centered."""
    iw, ih = Image.open(path).size
    box_ar = w / h; img_ar = iw / ih
    if img_ar > box_ar:
        dw = w; dh = w / img_ar
    else:
        dh = h; dw = h * img_ar
    px = x + (w - dw) / 2; py = y + (h - dh) / 2
    s.shapes.add_picture(str(path), Inches(px), Inches(py), Inches(dw), Inches(dh))


def caption(s, text, y=6.9):
    tf = box(s, 0.85, y, 11.6, 0.5)
    _set(tf.paragraphs[0], text, 13, MUTED, italic=True)


def table(s, rows, x, y, w, h, header_fill=ACCENT, font=13, col_widths=None,
          cell_colors=None):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for j, cw in enumerate(col_widths):
            gt.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.margin_left = Pt(6); c.margin_right = Pt(6)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            _set(p, str(val), font, RGBColor(0xFF,0xFF,0xFF) if i == 0 else INK,
                 bold=(i == 0 or j == 0), align=PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)
            if i == 0:
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                c.fill.solid(); c.fill.fore_color.rgb = BG if i % 2 else PANEL
            if cell_colors and (i, j) in cell_colors:
                c.text_frame.paragraphs[0].runs[0].font.color.rgb = cell_colors[(i, j)]
                c.text_frame.paragraphs[0].runs[0].font.bold = True
    return gt


# ════════════════════════════════════════════════════════════════════════════
# 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = slide()
panel = s.shapes.add_shape(1, Inches(0), Inches(0), SW, SH)
panel.fill.solid(); panel.fill.fore_color.rgb = INK; panel.line.fill.background()
tf = box(s, 0.9, 2.3, 11.5, 2.6)
_set(tf.paragraphs[0], "SYNTHETIC BANKING DATA", 15, RGBColor(0x9F,0xC0,0xE0), bold=True)
p = tf.add_paragraph(); _set(p, "A constrained, benchmarked generation pipeline", 40, BG, bold=True)
p = tf.add_paragraph()
_set(p, "5 methods (SDV + differential privacy) · constraints · LLM recommender", 20,
     RGBColor(0xC9,0xD3,0xDD))
tf = box(s, 0.9, 6.4, 11.5, 0.6)
_set(tf.paragraphs[0], "SDV · OpenDP SmartNoise · SDMetrics · DeepSeek  |  customers + transactions", 14,
     RGBColor(0x8A,0x97,0xA3))

# ════════════════════════════════════════════════════════════════════════════
# 2 — Problem
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Why synthetic data?", "The problem")
bullets(s, [
    "Real customer + transaction data is sensitive, regulated, and scarce.",
    "Teams still need realistic data to build & test models, demos, and pipelines.",
    "Naïve synthesis preserves single-column shapes but breaks what matters:",
    ("cross-table correlations — does income still predict the products bought?", 1),
    ("business rules — a product's category, a valid credit score, whole-number counts.", 1),
    "Goal: data that is statistically faithful AND provably valid — then prove it's useful.",
], gap=12)

# ════════════════════════════════════════════════════════════════════════════
# 3 — Pipeline architecture
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "End-to-end pipeline", "Architecture")

def stage(x, y, w, h, title, sub, color):
    sh = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = color
    tf = sh.text_frame; tf.word_wrap = True
    _set(tf.paragraphs[0], title, 16, BG, bold=True, align=PP_ALIGN.CENTER)
    p = tf.add_paragraph(); _set(p, sub, 11, RGBColor(0xEE,0xF2,0xF6), align=PP_ALIGN.CENTER)

def arrow(x, y):
    a = s.shapes.add_shape(13, Inches(x), Inches(y), Inches(0.5), Inches(0.5))
    a.fill.solid(); a.fill.fore_color.rgb = MUTED; a.line.fill.background()

yy, hh = 2.9, 1.7
stage(0.85, yy, 2.5, hh, "Seed data", "2,000 customers\n~24k transactions", INK)
arrow(3.45, yy+0.6)
stage(4.05, yy, 2.9, hh, "5 synthesizers", "SDV: HMA·CTGAN\nPAR·TVAE  +  DP", ACCENT)
arrow(7.05, yy+0.6)
stage(7.65, yy, 2.5, hh, "Evaluation", "SDMetrics +\ncustom metrics", METHOD["M3"])
arrow(10.25, yy+0.6)
stage(10.85, yy, 1.95, hh, "LLM", "DeepSeek\nrecommender", METHOD["M4"])
tf = box(s, 0.85, 4.9, 11.8, 1.5); tf.word_wrap = True
_set(tf.paragraphs[0], "Constraints (SDV CAG) are attached to every synthesizer before fitting —",
     17, INK, bold=True)
p = tf.add_paragraph()
_set(p, "so validity is guaranteed at generation time, not patched afterward.", 17, MUTED)

# ════════════════════════════════════════════════════════════════════════════
# 4 — Seed data
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "The ground truth: rule-driven seed data", "Seed data")
fit_image(s, REPORTS/"seed_overview.png", 0.8, 1.8, 12.0, 4.6)
caption(s, "Income↔occupation, credit eligibility, age→channel — the real structure all four methods must reproduce.")

# ════════════════════════════════════════════════════════════════════════════
# 5 — The four methods
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Five synthesis strategies — SDV vs OpenDP", "Methods")
rows = [
    ["", "Synthesizer", "Library", "Strength"],
    ["M1", "HMA + Gaussian Copula", "SDV", "Joint structure, native FK"],
    ["M2", "Independent CTGAN", "SDV", "Sharp marginals, timing"],
    ["M3", "CTGAN + PAR", "SDV", "Temporal / cross-table context"],
    ["M4", "Independent TVAE", "SDV", "Smooth, stable, correlated"],
    ["M5", "SmartNoise MST", "OpenDP", "(ε, δ)-differential privacy"],
]
cc = {(i+1, 0): METHOD[m] for i, m in enumerate(["M1","M2","M3","M4","M5"])}
table(s, rows, 0.85, 1.95, 11.6, 3.3, font=14,
      col_widths=[0.9, 3.9, 2.0, 4.8], cell_colors=cc)
bullets(s, ["M1–M4 (SDV) optimise statistical fidelity; M5 (SmartNoise/OpenDP) adds a formal privacy guarantee.",
            "Same seed data and metadata — only the synthesizer changes, so the comparison is clean."],
        y=5.5, size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════════
# 6 — Constraints
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Constraints make output valid, not just realistic", "SDV CAG constraints")
rows = [
    ["Field(s)", "Constraint", "Rule"],
    ["product_id ↔ product_category", "FixedCombinations", "category must match the product"],
    ["amount", "ScalarInequality ≥ 0", "no negative amounts"],
    ["num_dependents", "FixedIncrements(1)", "whole-number counts"],
    ["credit_score", "ScalarRange [300, 850]", "valid FICO range"],
    ["tenure_years ≤ age", "Inequality", "can't bank longer than you've lived"],
]
table(s, rows, 0.85, 2.0, 11.6, 3.0, font=14, col_widths=[4.0, 3.4, 4.2])
bullets(s, [
    "Attached via synth.add_constraints([...]) before fit — guaranteed in every sampled row.",
    "Scalar* use SDV's legacy dict form (not in sdv.cag); the rest are native CAG classes.",
], y=5.3, size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════════
# 7 — Constraint adherence
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Adherence: every rule holds (one caveat)", "Validity check")
rows = [
    ["Method", "prod/cat drift", "deps whole", "credit∈[300,850]", "tenure≤age", "amount≥0"],
    ["M1 HMA GC", "0", "✓", "✓", "✓", "✓"],
    ["M2 CTGAN", "0", "✓", "✓", "✓", "✓"],
    ["M3 CTGAN+PAR", "15", "✓", "✓", "✓", "✓"],
    ["M4 TVAE", "0", "✓", "✓", "✓", "✓"],
    ["M5 SmartNoise", "0", "✓", "✓", "✓", "✓"],
]
cc = {(1,1):GOOD,(2,1):GOOD,(4,1):GOOD,(5,1):GOOD,(3,1):WARN}
table(s, rows, 0.85, 2.0, 11.6, 3.0, font=13,
      col_widths=[2.6,2.2,1.6,2.2,1.5,1.5], cell_colors=cc)
bullets(s, [
    "M1 / M2 / M4 (SDV constraints) and M5 (DP + post-processing) satisfy every rule, 0 orphan FKs.",
    "M3 keeps 15 product/category drifts — PAR merges the two columns and rejects FixedCombinations.",
], y=5.3, size=14, gap=8)

# ════════════════════════════════════════════════════════════════════════════
# 8 — Evaluation metrics
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "How we measure quality", "Evaluation")
bullets(s, [
    "SDMetrics — standard fidelity:",
    ("Column shapes (marginals) · Column pair trends (within-table correlation)", 1),
    ("Diagnostic — referential integrity, boundaries, coverage", 1),
    "Custom metrics — what the use case actually needs:",
    ("Cross-table MAD ↓ — mean |Δ| of feature→product-category Spearman corr (real vs synth)", 1),
    ("Temporal realism — inter-arrival KS test + amount autocorrelation", 1),
    "Key idea: we score correlations and business signal, not just histograms.",
], gap=10, size=17)

# ── Metric definitions 1 — SDMetrics fidelity ───────────────────────────────
s = slide(); header(s, "Metrics — fidelity (SDMetrics)", "How they're computed")
rows = [
    ["Metric", "What it measures", "How it's computed", "↑/↓"],
    ["Cust. Shapes", "Per-column marginals (customers)",
     "avg KSComplement (num) / TVComplement (cat)", "↑"],
    ["Txn. Shapes", "Per-column marginals (transactions)",
     "same complements on transaction columns", "↑"],
    ["Cust. Pairs", "Pairwise column relationships",
     "avg CorrelationSimilarity / ContingencySimilarity", "↑"],
    ["Diagnostic", "Validity & structure (not fidelity)",
     "ranges · categories · keys · FK integrity", "↑"],
]
table(s, rows, 0.7, 2.05, 12.0, 3.0, font=13,
      col_widths=[1.9, 3.5, 5.0, 0.7])
bullets(s, [
    "KSComplement = 1 − KS distance · TVComplement = 1 − total-variation distance (1.0 = identical).",
    "Overall quality score = average of Column Shapes and Column Pair Trends.",
], y=5.3, size=14, gap=8)

# ── Metric definitions 2 — custom correlation & temporal ────────────────────
s = slide(); header(s, "Metrics — correlation, temporal & privacy (custom + SDMetrics)", "How they're computed")
rows = [
    ["Metric", "What it measures", "How it's computed", "↑/↓"],
    ["Cross-table MAD", "feature → product-category signal",
     "mean |Δ Spearman(feature, %category)|", "↓"],
    ["Within-table corr MAE", "customer column correlations",
     "mean |Δ Pearson| over off-diagonal pairs", "↓"],
    ["IA KS p-value", "transaction timing realism",
     "KS test on days-between-consecutive-txns", "↑"],
    ["Autocorr MAE", "sequential spend pattern",
     "|Δ mean lag-1 autocorrelation of amount|", "↓"],
    ["NewRowSynthesis", "verbatim copying (privacy)",
     "% synthetic rows not matching a real row", "↑"],
    ["DCR protection", "memorisation (privacy)",
     "synth→real nearest-neighbour dist vs baseline", "↑"],
]
table(s, rows, 0.7, 1.95, 12.0, 3.6, font=12,
      col_widths=[2.7, 3.3, 5.0, 0.7])
bullets(s, [
    "Autocorr = a column vs its own past (over time); corr MAE = between different columns (same row).",
    "Privacy metrics (SDMetrics) measure leakage; M5 alone adds a formal (ε, δ)-DP guarantee.",
], y=5.75, size=13, gap=7)

# ════════════════════════════════════════════════════════════════════════════
# 9 — Results dashboard
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "All methods, all metrics", "Results")
fit_image(s, REPORTS/"comparison_dashboard.png", 0.8, 1.7, 12.0, 4.7)
caption(s, "Normalised so higher = better on every axis. No single method dominates — the trade-offs are the story.")

# ── SDMetrics Quality Report showcase ────────────────────────────────────────
s = slide(); header(s, "SDMetrics Quality Report", "Fidelity scorecard")
fit_image(s, REPORTS/"quality_report.png", 0.7, 1.65, 9.0, 5.0)
tf = box(s, 9.9, 1.85, 3.1, 4.7); tf.word_wrap = True
_set(tf.paragraphs[0], "Overall = mean of\nColumn Shapes &\nColumn Pair Trends", 13, ACCENT, bold=True)
for t in ["M1 ≈ M4 lead overall (0.89) — strong on both halves.",
          "M5 wins customer pair trends (0.85) but on raw numeric correlation it's worst (see in-table slide).",
          "M2 best on transaction shapes (0.91), weak on pair trends.",
          "M3 trails everywhere (0.65) — PAR sacrifices marginals for sequence context."]:
    p = tf.add_paragraph(); _set(p, "• " + t, 12, INK); p.space_after = Pt(8)
caption(s, "SDMetrics QualityReport overall score decomposed into its sub-components, per method.")

# ════════════════════════════════════════════════════════════════════════════
# 10 — Cross-table heatmap
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Does income still predict products?", "Cross-table correlation")
fit_image(s, REPORTS/"cross_table_heatmap.png", 0.6, 1.8, 12.2, 4.5)
caption(s, "Feature→category Spearman correlations. Weak for every method — the cross-table signal is the hardest to keep.")

# Within-table correlation
s = slide(); header(s, "Within-customer correlation: M1 dominates", "In-table correlation")
fit_image(s, REPORTS/"in_table_correlation.png", 0.7, 1.65, 8.4, 4.9)
tf = box(s, 9.4, 1.9, 3.5, 4.4); tf.word_wrap = True
_set(tf.paragraphs[0], "MAE vs Real ↓", 15, ACCENT, bold=True)
for m, v, note in [("M1", "0.029", "Gaussian Copula"), ("M4", "0.056", ""),
                   ("M2", "0.132", ""), ("M3", "0.139", ""),
                   ("M5", "0.225", "DP cost")]:
    p = tf.add_paragraph()
    _set(p, f"{m}  {v}" + (f"   ({note})" if note else ""), 14,
         METHOD[m], bold=(m in ("M1", "M5"))); p.space_after = Pt(7)
p = tf.add_paragraph()
_set(p, "M1 reproduces income↔credit↔age ~2–8× better; M5 (best marginals) is worst on correlation.",
     12, MUTED, italic=True)
caption(s, "Pearson correlation of customer numeric columns — Real vs each method (MAE over off-diagonal pairs).")

# Correlation summary table (all methods, all correlation metrics)
s = slide(); header(s, "Correlation scorecard — all methods", "Correlation summary")
rows = [
    ["Method", "Within-table\ncorr MAE ↓", "Cross-table\nMAD ↓", "Cust. pair\ntrends ↑"],
    ["M1 HMA GC", "0.029", "0.379", "0.737"],
    ["M2 CTGAN", "0.132", "0.394", "0.567"],
    ["M3 CTGAN+PAR", "0.139", "0.296", "0.547"],
    ["M4 TVAE", "0.056", "0.392", "0.780"],
    ["M5 SmartNoise", "0.225", "0.387", "0.850"],
]
cc = {(i+1, 0): METHOD[m] for i, m in enumerate(["M1","M2","M3","M4","M5"])}
cc.update({(1,1): GOOD, (3,2): GOOD, (5,3): GOOD})   # winners: M1 within, M3 cross-table, M5 pairs
table(s, rows, 1.6, 2.1, 10.1, 3.2, font=14,
      col_widths=[3.1, 2.5, 2.3, 2.2], cell_colors=cc)
bullets(s, [
    "M1 owns the raw numeric correlation matrix (within-table MAE 0.029); M4 second (0.056).",
    "SDMetrics pair trends (num+cat) favours M5 (0.850), but its raw numeric correlations are worst.",
    "Cross-table (income→product) stays weak for all — M3 least bad.",
], y=5.55, size=13, gap=7)

# ── How differential privacy works ──────────────────────────────────────────
s = slide(); header(s, "How differential privacy works", "Method 5 · the guarantee", color=METHOD["M5"])
panel = s.shapes.add_shape(5, Inches(0.85), Inches(1.75), Inches(11.6), Inches(1.2))
panel.fill.solid(); panel.fill.fore_color.rgb = PANEL; panel.line.color.rgb = METHOD["M5"]
tf = panel.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0],
     "Add or remove any one person   →   P(output | with you)  ≤  e^ε · P(output | without you)",
     17, INK, bold=True, align=PP_ALIGN.CENTER)
p = tf.add_paragraph()
_set(p, "The two clouds of possible outputs stay nearly identical — your footprint hides in the noise.",
     13, MUTED, align=PP_ALIGN.CENTER)
bullets(s, [
    "A mathematical promise about the algorithm, not a property of the output — provable, not measured after the fact.",
    "ε (privacy budget): small ε = strong privacy, blurry data; large ε = sharp data, weak privacy. Pipeline ≈ 3 / table.",
    "δ: tiny probability the guarantee is allowed to fail (needed by the Gaussian mechanism).",
    "Mechanism: add calibrated noise to aggregates — noise ∝ sensitivity / ε (how much one person can move the result).",
    ("Never noise individual rows — noise the counts / histograms, then build data from those.", 1),
    "Composition: customers (ε=3) + transactions (ε=3) → ε ≈ 6 for someone in both tables.",
    "Post-processing immunity: anything computed from DP output stays DP — you can't 'use it up'.",
], y=3.15, size=14.5, gap=8)

# ── How SmartNoise / MST builds M5 ───────────────────────────────────────────
s = slide(); header(s, "How SmartNoise / MST builds M5", "Method 5 · the mechanism", color=METHOD["M5"])
bullets(s, [
    "MST (private-PGM family): marginal-based — it learns noisy summaries and never memorises rows.",
    "1 · Discretize — bucket every column into bins; privately learns the ranges (preprocessor_eps).",
    "2 · Noisy 1-way marginals — per-column histograms with Gaussian noise added to every count.",
    "3 · Maximum spanning tree — privately find the most-correlated column pairs and keep the tree that connects all columns; only those 2-way marginals are measured (+ noise).",
    "4 · Fit a graphical model (Markov Random Field) consistent with the noisy marginals — pure post-processing, zero extra ε.",
    "5 · Sample new rows from the fitted model.",
    "Runs once per table (customers, transactions) → per-table ε composes to ≈ 6.",
], y=1.95, size=15, gap=8)
panel = s.shapes.add_shape(5, Inches(0.85), Inches(5.7), Inches(11.6), Inches(1.25))
panel.fill.solid(); panel.fill.fore_color.rgb = PANEL; panel.line.color.rgb = METHOD["M5"]
tf = panel.text_frame; tf.word_wrap = True
_set(tf.paragraphs[0], "Why M5 looks the way it does", 13, METHOD["M5"], bold=True)
p = tf.add_paragraph()
_set(p, "Guarantee by construction · only measured marginals survive (off-tree structure like per-customer "
        "sequencing is lost → autocorr collapses) · faithful marginals sit close to real rows → low DCR.",
     12.5, INK)

# Privacy — guarantee vs measured
s = slide(); header(s, "Privacy: guarantee ≠ measured protection", "Privacy", color=METHOD["M5"])
fit_image(s, REPORTS/"privacy_dcr.png", 0.7, 1.75, 7.6, 4.5)
tf = box(s, 8.5, 1.9, 4.4, 4.6); tf.word_wrap = True
_set(tf.paragraphs[0], "Two notions of privacy", 16, METHOD["M5"], bold=True)
for t in ["M5 (SmartNoise) is the ONLY method with a formal (ε≈6, δ)-DP guarantee.",
          "But on measured DCR it is the LEAST protected (0.57) — M3 most (0.74).",
          "DP bounds one person's worst-case influence; DCR measures average closeness to real rows.",
          "M5's faithful marginals sit closest to real data → lowest DCR.",
          "NewRowSynthesis = 1.0 for all (no verbatim copies).",
          "Fidelity and record-distance trade off — a guarantee ≠ best DCR."]:
    p = tf.add_paragraph(); _set(p, "• " + t, 12.5, INK); p.space_after = Pt(7)
caption(s, "DCRBaselineProtection on the customers table (higher = more protected) — measured with SDMetrics.")

# ── Measurable privacy: membership-inference audit ───────────────────────────
s = slide(); header(s, "Can we measure privacy? A membership-inference audit", "Measured privacy", color=METHOD["M5"])
fit_image(s, REPORTS/"privacy_audit.png", 0.45, 1.7, 9.05, 4.35)
tf = box(s, 9.65, 1.85, 3.35, 4.7); tf.word_wrap = True
_set(tf.paragraphs[0], "ε is accounted, not measured", 13.5, METHOD["M5"], bold=True)
for t in ["Guarantee (ε) = worst-case bookkeeping from the accountant — can't be read off the data.",
          "Measurable instead: attacker AUC separating training members from fresh non-members.",
          "B=300 bootstrap: every method's AUC interval straddles 0.50 → no significant leakage.",
          "Clopper-Pearson empirical ε = 0 for all — no certifiable leak, far under M5's ε≈6 budget.",
          "Average AUC can miss worst-case records; M5's DP bound is the only per-record, any-scale guarantee."]:
    p = tf.add_paragraph(); _set(p, "• " + t, 11, INK); p.space_after = Pt(6)
caption(s, "Members = real training rows; non-members = fresh draw from the same process. AUC 0.5 = no leakage; ε via Clopper-Pearson 95% lower bound (B=300 bootstrap CIs on AUC).")

# ════════════════════════════════════════════════════════════════════════════
# 11 — Distributions
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Marginal fidelity", "Distributions")
# four separate charts in a 2x2 grid — easier to read than one 4-in-1 strip
gx, gy, gw, gh, pad = 0.85, 1.75, 11.6, 4.75, 0.25
cw, ch = (gw - pad) / 2, (gh - pad) / 2
for k, img in enumerate(["dist_age.png", "dist_income.png", "dist_credit.png", "dist_tenure.png"]):
    r, c = divmod(k, 2)
    fit_image(s, REPORTS/img, gx + c*(cw+pad), gy + r*(ch+pad), cw, ch)
caption(s, "Real (grey fill) vs each method (coloured outline) — per variable for easier comparison.")

# ════════════════════════════════════════════════════════════════════════════
# 12 — Verdict
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Results at 2,000 seeds · ~12 txns — SDV vs DP", "Verdict")
rows = [
    ["Metric", "M1", "M2", "M3", "M4", "M5", "Winner"],
    ["Overall quality", "0.888", "0.854", "0.647", "0.887", "0.861", "M1≈M4"],
    ["Diagnostic / FK", "1.000", "1.000", "0.744", "1.000", "1.000", "tie"],
    ["Cust. pair trends", "0.737", "0.567", "0.547", "0.780", "0.850", "M5"],
    ["Txn. column shapes", "0.865", "0.908", "0.783", "0.826", "0.856", "M2"],
    ["Cross-table MAD ↓", "0.379", "0.394", "0.296", "0.392", "0.387", "M3"],
    ["Autocorr MAE ↓", "0.079", "0.011", "0.041", "0.028", "0.569", "M2"],
    ["Diff. privacy", "✗", "✗", "✗", "✗", "✓", "M5"],
    ["DCR protection ↑", "0.61", "0.69", "0.74", "0.65", "0.57", "M3"],
]
table(s, rows, 0.7, 1.9, 9.2, 4.2, font=12,
      col_widths=[2.5,0.95,0.95,0.95,0.95,0.95,1.95])
tf = box(s, 10.15, 1.9, 2.9, 4.3); tf.word_wrap = True
_set(tf.paragraphs[0], "Pick by need", 16, ACCENT, bold=True)
for t in ["M1 ≈ M4 — best all-round fidelity",
          "M1 — demographic correlations (recommender)",
          "M2 — transaction shapes & timing",
          "M5 — the only privacy guarantee",
          "M3 — cross-table & DCR, low utility"]:
    p = tf.add_paragraph(); _set(p, "• " + t, 13, INK); p.space_after = Pt(7)
p = tf.add_paragraph()
_set(p, "More data closed M5's quality lead; 12 txns/cust exposed its temporal gap.",
     12, METHOD["M5"], italic=True)

# ── Bootstrap confidence intervals ───────────────────────────────────────────
s = slide(); header(s, "Are the differences real? 95% bootstrap CIs", "Robustness")
fit_image(s, REPORTS/"bootstrap_ci.png", 0.5, 1.75, 8.8, 4.9)
tf = box(s, 9.55, 1.9, 3.45, 4.8); tf.word_wrap = True
_set(tf.paragraphs[0], "What the overlap tells us", 14, ACCENT, bold=True)
for t in ["B=200 resamples of customers (+ their txns); whiskers = 2.5–97.5%.",
          "M1 ≈ M4 on quality — CIs overlap, so the tie is real, not noise.",
          "Cross-table: only M3 separates; M1/M2/M4/M5 are statistically tied.",
          "M5 leads pair trends (lower CI 0.80 > M4's 0.79) — a genuine edge.",
          "M5's autocorr (~0.57) and lowest DCR are clear-cut, non-overlapping.",
          "M5's quality CI is the widest — DP randomness makes it least stable."]:
    p = tf.add_paragraph(); _set(p, "• " + t, 11.5, INK); p.space_after = Pt(7)
caption(s, "Median ± 95% percentile CI per metric (privacy at B=40, subsample 500). Non-overlapping intervals = a real difference.")

# ════════════════════════════════════════════════════════════════════════════
# 13 — LLM application
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Closing the loop: LLM recommendations", "Application", color=METHOD["M4"])
tf = box(s, 0.85, 1.85, 11.6, 0.9); tf.word_wrap = True
_set(tf.paragraphs[0], "Feed a synthetic M1 customer profile to DeepSeek (deepseek-chat) → 3 ranked products.",
     17, INK)
p = tf.add_paragraph()
_set(p, "Synthetic customer C01502 — 39, Retired, income ~$135k, credit 699", 14, MUTED, italic=True)
rows = [
    ["#", "Product", "Why"],
    ["1", "Credit Card Basic", "solid credit, high income, holds no credit card → build history"],
    ["2", "Investment Fund A", "high income + long tenure → diversify, low entry, medium risk"],
    ["3", "Health Insurance Basic", "already holds premium health → low-cost secondary coverage"],
]
table(s, rows, 0.85, 3.0, 11.6, 2.2, font=14, header_fill=METHOD["M4"],
      col_widths=[0.7, 3.6, 7.3])
bullets(s, ["Catalog cached in the system prompt — repeat calls hit 512 cached tokens.",
            "Proves the synthetic data is fit for a real downstream task."],
        y=5.4, size=15, gap=8)

# ════════════════════════════════════════════════════════════════════════════
# 14 — Takeaways
# ════════════════════════════════════════════════════════════════════════════
s = slide(); header(s, "Takeaways", "Wrap-up")
bullets(s, [
    "Method choice is task-dependent — and the leaderboard moved when we scaled to 2k × 12 txns.",
    ("Privacy guarantee? M5. All-round fidelity / correlations? M1 ≈ M4. Transaction timing? M2.", 1),
    "More data closed M5's quality lead; more transactions exposed its missing temporal model.",
    "A DP guarantee (M5) ≠ best measured DCR (M3): guarantee and empirical leakage are different axes.",
    "Constraints turn 'realistic-looking' data into provably valid data.",
    "End-to-end & reproducible: seed → synthesize → evaluate (+privacy) → recommend, one notebook.",
], gap=11, size=17)
tf = box(s, 0.85, 6.5, 11.6, 0.6)
_set(tf.paragraphs[0], "Code: src/{methods,constraints,evaluate,llm_suggest}.py  ·  synthetic_data_pipeline.ipynb",
     13, MUTED, italic=True)

out = Path("slides.pptx")
prs.save(out)
print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
